#!/usr/bin/env python3
"""Generate role-specific migration packs from an analysis folder."""

from __future__ import annotations

import argparse
import csv
import json
import re
from pathlib import Path


PERSONAS = {
    "ceo": {
        "title": "CEO Board Pack",
        "files": ["persona-ceo-summary.md", "board-ceo-narrative.md", "ai-value-tracker.md", "ai-cost-model.md", "ai-risk-heatmap.md", "steering-committee-pack.md"],
        "score": ["ai-value-tracker.md", "ai-cost-model.md", "persona-ceo-summary.md"],
    },
    "cio": {
        "title": "CIO Architecture Pack",
        "files": ["persona-cio-architecture-view.md", "ai-before-after-architecture.md", "ai-upgrade-path-decision.md", "ai-dependency-graph.md", "ai-adrs.md"],
        "score": ["persona-cio-architecture-view.md", "ai-upgrade-path-decision.md", "ai-dependency-graph.md"],
    },
    "ciso": {
        "title": "CISO Security Gate Pack",
        "files": ["persona-ciso-security-view.md", "ciso-security-gate-pack.md", "ai-quality-gates.md", "ai-risk-mitigation-playbooks.md", "ai-evidence-confidence.md"],
        "score": ["persona-ciso-security-view.md", "ciso-security-gate-pack.md", "ai-quality-gates.md"],
    },
    "pm": {
        "title": "Project Manager Control Pack",
        "files": ["persona-project-manager-control-view.md", "raid-log.md", "raci-matrix.md", "weekly-status-report.md", "steering-committee-pack.md", "project-operating-model.md", "ai-wave-roadmap.md"],
        "score": ["persona-project-manager-control-view.md", "raid-log.md", "weekly-status-report.md"],
    },
    "team": {
        "title": "Team Execution Pack",
        "files": ["persona-team-member-task-view.md", "team-execution-pack.md", "role-based-prompt-library.md", "project-onboarding-guide.md", "ai-migration-backlog.md", "ai-workshop-questions.md"],
        "score": ["persona-team-member-task-view.md", "team-execution-pack.md", "ai-migration-backlog.md"],
    },
}


def read(path: Path) -> str:
    return path.read_text(encoding="utf-8") if path.exists() else ""


def table_rows(text: str) -> int:
    return sum(1 for line in text.splitlines() if line.startswith("|") and "---" not in line)


def readiness_score(analysis_dir: Path, persona: str) -> tuple[int, list[str]]:
    cfg = PERSONAS[persona]
    present = [name for name in cfg["score"] if (analysis_dir / name).exists()]
    text = "\n".join(read(analysis_dir / name).lower() for name in cfg["score"])
    risk_hits = len(re.findall(r"\b(high|critical|blocked|direct-sql|overlayering|security|missing)\b", text))
    evidence_hits = len(re.findall(r"\b(confidence|owner|evidence|decision|gate|mitigation)\b", text))
    base = int((len(present) / max(len(cfg["score"]), 1)) * 55)
    evidence = min(evidence_hits * 2, 25)
    risk_penalty = min(risk_hits, 30)
    score = max(5, min(100, base + evidence + 20 - risk_penalty))
    actions = []
    if score < 50:
        actions.append("Immediate evidence review and leadership decision workshop required.")
    if risk_hits > 10:
        actions.append("Reduce high-risk scope before committing delivery baseline.")
    if len(present) < len(cfg["score"]):
        actions.append("Regenerate analysis pack because some persona source files are missing.")
    if not actions:
        actions.append("Proceed with gate review and maintain weekly evidence updates.")
    return score, actions


def build_markdown(analysis_dir: Path, persona: str) -> str:
    cfg = PERSONAS[persona]
    score, actions = readiness_score(analysis_dir, persona)
    parts = [
        f"# {cfg['title']}",
        "",
        "## Readiness Score",
        "",
        f"| Persona | Score | Interpretation |",
        "| --- | --- | --- |",
        f"| {persona.upper()} | {score}/100 | {'Ready' if score >= 75 else 'Needs control' if score >= 50 else 'At risk'} |",
        "",
        "## Next Actions",
        "",
        *[f"- {action}" for action in actions],
        "",
        "## Included Evidence",
        "",
    ]
    for name in cfg["files"]:
        path = analysis_dir / name
        if path.exists():
            text = read(path)
            parts.extend([f"### {name}", "", text[:3500].rstrip(), ""])
        else:
            parts.extend([f"### {name}", "", "Missing from analysis folder. Regenerate or validate source analysis.", ""])
    return "\n".join(parts).rstrip() + "\n"


def export_xlsx(scores: dict[str, int], output: Path) -> Path | None:
    try:
        from openpyxl import Workbook
    except ImportError:
        return None
    wb = Workbook()
    ws = wb.active
    ws.title = "Readiness"
    ws.append(["Persona", "Score", "Status"])
    for persona, score in scores.items():
        ws.append([persona.upper(), score, "Ready" if score >= 75 else "Needs control" if score >= 50 else "At risk"])
    path = output / "persona-readiness.xlsx"
    wb.save(path)
    return path


def export_pptx(scores: dict[str, int], output: Path) -> Path | None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return None
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AX to D365FO Persona Readiness"
    slide.placeholders[1].text = "CEO, CIO, CISO, PM, and Team view"
    for persona, score in scores.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"{persona.upper()} Readiness: {score}/100"
        frame = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(4)).text_frame
        p = frame.paragraphs[0]
        p.text = "Ready" if score >= 75 else "Needs control" if score >= 50 else "At risk"
        p.font.size = Pt(28)
    path = output / "persona-readiness-deck.pptx"
    prs.save(path)
    return path


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("analysis_dir")
    parser.add_argument("--persona", choices=[*PERSONAS.keys(), "all"], default="all")
    parser.add_argument("--output", default="persona-packs")
    parser.add_argument("--office", action="store_true", help="Also generate XLSX/PPTX readiness exports.")
    args = parser.parse_args()

    analysis_dir = Path(args.analysis_dir)
    output = Path(args.output)
    output.mkdir(parents=True, exist_ok=True)
    personas = list(PERSONAS) if args.persona == "all" else [args.persona]
    scores = {}
    for persona in personas:
        score, _ = readiness_score(analysis_dir, persona)
        scores[persona] = score
        (output / f"{persona}-pack.md").write_text(build_markdown(analysis_dir, persona), encoding="utf-8")
    (output / "readiness-scores.json").write_text(json.dumps(scores, indent=2), encoding="utf-8")
    with (output / "readiness-scores.csv").open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerow(["Persona", "Score"])
        writer.writerows([[k, v] for k, v in scores.items()])
    if args.office:
        export_xlsx(scores, output)
        export_pptx(scores, output)
    print(f"Generated persona pack into {output.resolve()}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
