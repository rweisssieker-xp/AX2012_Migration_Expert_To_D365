#!/usr/bin/env python3
"""Export analyzer outputs to XLSX and PPTX deliverables."""

from __future__ import annotations

import argparse
import csv
from pathlib import Path


def read_csv(path: Path) -> list[list[str]]:
    with path.open("r", encoding="utf-8-sig", newline="") as handle:
        return list(csv.reader(handle))


def read_md_table(path: Path) -> list[list[str]]:
    rows = []
    for line in path.read_text(encoding="utf-8").splitlines():
        if line.startswith("|") and "---" not in line:
            rows.append([cell.strip() for cell in line.strip("|").split("|")])
    return rows


def export_xlsx(analysis_dir: Path, output: Path) -> Path:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill
        from openpyxl.utils import get_column_letter
    except ImportError as exc:
        raise SystemExit("openpyxl is required for XLSX export.") from exc

    wb = Workbook()
    wb.remove(wb.active)
    for source in sorted(analysis_dir.glob("ai-*.md")):
        rows = read_md_table(source)
        if not rows:
            continue
        ws = wb.create_sheet(source.stem[:31])
        for row in rows:
            ws.append(row)
        for cell in ws[1]:
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="0078D4")
        ws.freeze_panes = "A2"
        for col in range(1, ws.max_column + 1):
            width = min(max(len(str(ws.cell(row=row, column=col).value or "")) for row in range(1, ws.max_row + 1)) + 2, 60)
            ws.column_dimensions[get_column_letter(col)].width = width
    output.mkdir(parents=True, exist_ok=True)
    path = output / "migration-analysis.xlsx"
    wb.save(path)
    return path


def export_control_workbooks(analysis_dir: Path, output: Path) -> list[Path]:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill
    except ImportError as exc:
        raise SystemExit("openpyxl is required for XLSX export.") from exc

    specs = {
        "pmo-control-workbook.xlsx": [
            ("Workstream", "Status", "Owner", "Risk", "Next Action"),
            ("Finance", "Needs control", "PMO", "Reconciliation", "Close sign-off evidence"),
            ("Data", "Needs control", "Data Lead", "Quality", "Run data profiling"),
            ("Cutover", "Blocked", "Cutover Lead", "Rehearsal", "Complete rehearsal scorecard"),
        ],
        "commerce-readiness-workbook.xlsx": [
            ("Area", "Score", "Status", "Evidence", "Next Action"),
            ("CSU", "75", "Ready", "Readiness report", "Validate performance baseline"),
            ("POS Offline", "60", "Needs control", "Offline runbook", "Execute offline recovery test"),
            ("Payments/PCI", "50", "Blocked", "PCI gate", "Obtain security approval"),
        ],
        "evidence-vault-workbook.xlsx": [
            ("Evidence", "Owner", "Required", "Present", "Gate"),
            ("CISO approval", "CISO", "Yes", "No", "Go-live"),
            ("Finance reconciliation", "CFO/Finance", "Yes", "No", "Go-live"),
            ("UAT sign-off", "Business", "Yes", "No", "Go-live"),
        ],
    }
    output.mkdir(parents=True, exist_ok=True)
    paths = []
    for filename, rows in specs.items():
        wb = Workbook()
        ws = wb.active
        ws.title = "Control"
        for row in rows:
            ws.append(row)
        for cell in ws[1]:
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="1F4E79")
        ws.freeze_panes = "A2"
        ws2 = wb.create_sheet("Source")
        ws2.append(["Analysis folder", str(analysis_dir.resolve())])
        path = output / filename
        wb.save(path)
        paths.append(path)
    return paths


def export_pptx(analysis_dir: Path, output: Path) -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise SystemExit("python-pptx is required for PPTX export.") from exc

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "AX to D365FO Migration Command Center"
    title_slide.placeholders[1].text = "Executive summary generated from inventory analysis"

    summary = read_md_table(analysis_dir / "ai-analysis-summary.md")
    deck_sections = [
        ("Complexity & Scope", summary[:8]),
        ("Top Risks", read_md_table(analysis_dir / "ai-risk-radar.md")[:8]),
        ("Wave Roadmap", read_md_table(analysis_dir / "ai-wave-roadmap.md")[:8]),
        ("Quality Gates", read_md_table(analysis_dir / "ai-quality-gates.md")[:8]),
        ("Cost Model", read_md_table(analysis_dir / "ai-cost-model.md")[:8]),
    ]
    for title, rows in deck_sections:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        tx = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.8), Inches(4.8)).text_frame
        tx.clear()
        for row in rows[1:]:
            p = tx.add_paragraph()
            p.text = " | ".join(row[:4])
            p.font.size = Pt(16)
    output.mkdir(parents=True, exist_ok=True)
    path = output / "migration-executive-deck.pptx"
    prs.save(path)
    return path


def export_role_decks(analysis_dir: Path, output: Path) -> list[Path]:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise SystemExit("python-pptx is required for PPTX export.") from exc

    decks = {
        "ceo-migration-value-deck.pptx": [
            ("CEO Migration Value", "Value, risk, confidence, benefits realization"),
            ("Board Signals", "Go-live confidence | budget exposure | business disruption | value leakage"),
            ("Decision Ask", "Approve scope controls, external sign-offs, and value tracking cadence"),
        ],
        "cio-architecture-deck.pptx": [
            ("CIO Architecture", "Target architecture, integrations, ALM, data and platform readiness"),
            ("Architecture Risks", "Unsupported customizations | integration resilience | data entity gaps"),
            ("Modernization Path", "Retire legacy patterns, add observability, plan Fabric/Power BI roadmap"),
        ],
        "ciso-gate-deck.pptx": [
            ("CISO Gate", "Security evidence, roles, secrets, PCI, privileged access and audit posture"),
            ("Blockers", "Missing CISO approval | payment evidence | service account register | attack surface map"),
            ("Approval Boundary", "External security/PCI sign-off remains a mandatory independent gate"),
        ],
        "board-risk-forecast-deck.pptx": [
            ("Board Risk Forecast", "Executive risk model based on evidence, testing, cutover and scope controls"),
            ("Forecast Drivers", "Evidence gaps | late scope | failed rehearsals | reconciliation defects"),
            ("Required Control", "Weekly board risk review until all critical gates are ready"),
        ],
    }
    output.mkdir(parents=True, exist_ok=True)
    paths = []
    for filename, slides in decks.items():
        prs = Presentation()
        for idx, (title, body) in enumerate(slides):
            slide = prs.slides.add_slide(prs.slide_layouts[0] if idx == 0 else prs.slide_layouts[5])
            slide.shapes.title.text = title
            box = slide.placeholders[1] if idx == 0 else slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.7), Inches(4.5))
            frame = box.text_frame
            frame.clear()
            for part in body.split(" | "):
                p = frame.add_paragraph()
                p.text = part
                p.font.size = Pt(18)
        path = output / filename
        prs.save(path)
        paths.append(path)
    return paths


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("analysis_dir")
    parser.add_argument("--output", default="migration-exports")
    args = parser.parse_args()

    analysis_dir = Path(args.analysis_dir)
    output = Path(args.output)
    exports = [export_xlsx(analysis_dir, output), export_pptx(analysis_dir, output)]
    exports.extend(export_control_workbooks(analysis_dir, output))
    exports.extend(export_role_decks(analysis_dir, output))
    for path in exports:
        print(f"Exported {path.resolve()}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
