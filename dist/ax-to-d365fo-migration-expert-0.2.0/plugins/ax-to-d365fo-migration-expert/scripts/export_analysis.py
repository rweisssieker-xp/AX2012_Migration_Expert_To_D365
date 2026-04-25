#!/usr/bin/env python3
"""Export analyzer outputs to XLSX and PPTX deliverables."""

from __future__ import annotations

import argparse
import csv
from pathlib import Path


def read_csv(path: Path) -> list[list[str]]:
    with path.open("r", encoding="utf-8-sig", newline="") as handle:
        return list(csv.reader(handle))


def read_md_table(path: Path) -> list[list[str]]:
    rows = []
    for line in path.read_text(encoding="utf-8").splitlines():
        if line.startswith("|") and "---" not in line:
            rows.append([cell.strip() for cell in line.strip("|").split("|")])
    return rows


def export_xlsx(analysis_dir: Path, output: Path) -> Path:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill
        from openpyxl.utils import get_column_letter
    except ImportError as exc:
        raise SystemExit("openpyxl is required for XLSX export.") from exc

    wb = Workbook()
    wb.remove(wb.active)
    for source in sorted(analysis_dir.glob("ai-*.md")):
        rows = read_md_table(source)
        if not rows:
            continue
        ws = wb.create_sheet(source.stem[:31])
        for row in rows:
            ws.append(row)
        for cell in ws[1]:
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="0078D4")
        ws.freeze_panes = "A2"
        for col in range(1, ws.max_column + 1):
            width = min(max(len(str(ws.cell(row=row, column=col).value or "")) for row in range(1, ws.max_row + 1)) + 2, 60)
            ws.column_dimensions[get_column_letter(col)].width = width
    output.mkdir(parents=True, exist_ok=True)
    path = output / "migration-analysis.xlsx"
    wb.save(path)
    return path


def export_pptx(analysis_dir: Path, output: Path) -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise SystemExit("python-pptx is required for PPTX export.") from exc

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "AX to D365FO Migration Command Center"
    title_slide.placeholders[1].text = "Executive summary generated from inventory analysis"

    summary = read_md_table(analysis_dir / "ai-analysis-summary.md")
    deck_sections = [
        ("Complexity & Scope", summary[:8]),
        ("Top Risks", read_md_table(analysis_dir / "ai-risk-radar.md")[:8]),
        ("Wave Roadmap", read_md_table(analysis_dir / "ai-wave-roadmap.md")[:8]),
        ("Quality Gates", read_md_table(analysis_dir / "ai-quality-gates.md")[:8]),
        ("Cost Model", read_md_table(analysis_dir / "ai-cost-model.md")[:8]),
    ]
    for title, rows in deck_sections:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        tx = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.8), Inches(4.8)).text_frame
        tx.clear()
        for row in rows[1:]:
            p = tx.add_paragraph()
            p.text = " | ".join(row[:4])
            p.font.size = Pt(16)
    output.mkdir(parents=True, exist_ok=True)
    path = output / "migration-executive-deck.pptx"
    prs.save(path)
    return path


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("analysis_dir")
    parser.add_argument("--output", default="migration-exports")
    args = parser.parse_args()

    analysis_dir = Path(args.analysis_dir)
    output = Path(args.output)
    xlsx = export_xlsx(analysis_dir, output)
    pptx = export_pptx(analysis_dir, output)
    print(f"Exported {xlsx.resolve()}")
    print(f"Exported {pptx.resolve()}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
